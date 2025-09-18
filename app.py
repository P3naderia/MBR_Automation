# -*- coding: utf-8 -*-
# python -m streamlit run test.py
# -*- coding: utf-8 -*-
# -*- coding: utf-8 -*-
import os, re, io, tempfile, shutil
from datetime import datetime
from typing import Dict, List, Optional, Tuple

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import streamlit as st
import urllib.request

# python-pptx
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR  # ← 색상 안전 처리용



def download_and_setup_font():
    # 나눔고딕 폰트 다운로드
    font_url = "https://github.com/naver/nanumfont/raw/master/fonts/NanumFontSetup_TTF_GOTHIC/NanumGothic.ttf"
    font_path = "NanumGothic.ttf"
    
    # 폰트 파일이 없으면 다운로드
    if not os.path.exists(font_path):
        urllib.request.urlretrieve(font_url, font_path)
    
    # 폰트 등록
    fm.fontManager.addfont(font_path)
    font_prop = fm.FontProperties(fname=font_path)
    
    # matplotlib에 폰트 설정
    plt.rc('font', family=font_prop.get_name())
    plt.rcParams['axes.unicode_minus'] = False
    
    return font_prop.get_name()

# 폰트 설정 실행
try:
    font_name = download_and_setup_font()
    print(f"Font loaded: {font_name}")
except Exception as e:
    print(f"Font loading failed: {e}")
    plt.rc('font', family='DejaVu Sans')
# =========================
# Global style / constants
# =========================
# 상단에 FONT_PATH 정의 추가


MONTH_LABELS = ['1월','2월','3월','4월','5월','6월','7월','8월','9월','10월','11월','12월']
PALETTE = {
    "primary": "#2F80ED", "green": "#27AE60", "orange": "#F2994A",
    "purple": "#9B51E0", "red": "#EB5757", "gray": "#BDBDBD", "dark": "#4F4F4F",
    "sp_fill": "#A9CEF8", "sb_fill": "#FAD4AD", "sd_fill": "#DCC4F6", "ba_fill": "#BFD3F2",
}
TOPCAT_METRICS = {}   # 그래프11에서 채움
TOPASIN_METRICS = {}  # Top1/Top2 ASIN 텍스트 마커용  ← 추가
GRAPH_ROOT: Optional[str] = None  # 세션 임시 작업폴더

# =========================
# Streamlit 기본 설정
# =========================
st.set_page_config(page_title="MBR PPT 자동 생성기 (3파일 업로드)", layout="wide")
st.title("📊 MBR PowerPoint 자동 생성기")
st.markdown("""
1) **CID 레벨 CSV/XLSX**, 2) **ASIN 레벨 CSV/XLSX**, 3) **PPT 템플릿(.pptx)** 을 업로드하세요.  
버튼을 누르면 **그래프(1~12, 16~17, 18~24) 생성 → 템플릿의 '그래프n' 자리 자동 삽입 → 텍스트/표 마커 치환 + YTD 테이블 채움 → 완성 PPT 다운로드**까지 자동 처리합니다.
""")

# =========================
# 파일 업로드
# =========================
cid_up = st.file_uploader("📁 CID 레벨 데이터 (CSV 또는 XLSX)", type=["csv","xlsx"], key="cid")
asin_up = st.file_uploader("📁 ASIN 레벨 데이터 (CSV 또는 XLSX)", type=["csv","xlsx"], key="asin")
ppt_up = st.file_uploader("📄 PowerPoint 템플릿 (.pptx)", type=["pptx"], key="ppt")

# =========================
# 공통 유틸
# =========================
def _set_korean_font_if_possible():
    try:
        if os.path.exists(FONT_PATH):
            font_prop = fm.FontProperties(fname=FONT_PATH)
            plt.rcParams['font.family'] = font_prop.get_name()
    except:
        pass

def ensure_graphs_folder() -> str:
    base = GRAPH_ROOT if GRAPH_ROOT else os.getcwd()
    graphs_folder = os.path.join(base, "graphs")
    os.makedirs(graphs_folder, exist_ok=True)
    return graphs_folder

def _save_fig(fig, name):
    graphs = ensure_graphs_folder()
    path = os.path.join(graphs, f"{name}.png")
    fig.savefig(path, dpi=300, bbox_inches="tight")
    plt.close(fig)
    return path

def parse_date_any(x):
    if pd.isna(x): return None
    try: return pd.to_datetime(str(x)).to_pydatetime().replace(day=1)
    except: return None

def parse_number_any(x, pct_to_100=False):
    if pd.isna(x): return None
    s = str(x).replace(',', '').replace('$', '').strip()
    m = re.search(r'[-+]?\d*\.?\d+', s)
    if not m: return None
    val = float(m.group(0))
    if '%' in s: return val
    return val*100 if pct_to_100 else val

def finalize_year_month(df, year_col='year', month_col='month'):
    df[year_col]  = pd.to_numeric(df[year_col], errors='coerce')
    df[month_col] = pd.to_numeric(df[month_col], errors='coerce')
    df = df.dropna(subset=[year_col, month_col]).copy()
    df[year_col]  = df[year_col].astype(int)
    df[month_col] = df[month_col].astype(int)
    df = df.sort_values([year_col, month_col]).reset_index(drop=True)
    df['date_str'] = df[year_col].astype(str) + '-' + df[month_col].astype(str).str.zfill(2)
    return df

def monthly_agg(dates, values, agg='mean'):
    rows=[]
    for d,v in zip(dates, values):
        dt = parse_date_any(d)
        val = parse_number_any(v) if v is not None else None
        if dt is None or val is None: continue
        rows.append({'year': dt.year, 'month': dt.month, 'value': val})
    if not rows: return pd.DataFrame(columns=['year','month','value','date_str'])
    t = pd.DataFrame(rows)
    gp = getattr(t.groupby(['year','month'])['value'], agg)()
    out = gp.reset_index()
    return finalize_year_month(out, 'year', 'month')

def _bi_theme(ax, ygrid=True):
    ax.set_facecolor("white")
    if ax.figure is not None: ax.figure.set_facecolor("white")
    for side in ["top","right"]: ax.spines[side].set_visible(False)
    for side in ["left","bottom"]: ax.spines[side].set_color("#BDBDBD")
    ax.tick_params(colors=PALETTE["dark"], labelsize=10)
    if ygrid: ax.yaxis.grid(True, color="#E6E6E6", linestyle="-", linewidth=1)
    ax.xaxis.grid(False)

def _yfmt_decimal(dec=1, suffix=""):
    return plt.FuncFormatter(lambda x, pos: f"{x:.{dec}f}{suffix}")

def _yfmt_k(dec=1):
    return plt.FuncFormatter(lambda x, pos: f"{x/1000:.{dec}f}K")

def _label_last(ax, xs, ys, text, dy=6):
    if len(xs)==0: return
    ax.annotate(text, (xs[-1], ys[-1]), textcoords="offset points", xytext=(0,dy),
                ha="center", va="bottom", fontsize=9,
                bbox=dict(boxstyle="round,pad=0.25", fc="white", ec="#DDDDDD", alpha=0.9))

# =========================
# 색상 안전 스냅샷/적용 (테마색 대응)
# =========================
def _snapshot_color(colorfmt):
    """python-pptx ColorFormat을 안전하게 스냅샷 (RGB 또는 테마)"""
    snap = {"type": None, "rgb": None, "theme": None}
    try:
        if colorfmt is None or colorfmt.type is None:
            return snap
        if colorfmt.type == MSO_COLOR_TYPE.RGB and colorfmt.rgb is not None:
            snap["type"] = "rgb"
            snap["rgb"] = colorfmt.rgb
            return snap
        if colorfmt.type == MSO_COLOR_TYPE.SCHEME and colorfmt.theme_color is not None:
            snap["type"] = "theme"
            snap["theme"] = colorfmt.theme_color
            return snap
    except Exception:
        pass
    return snap

def _apply_color(run, snap):
    """스냅샷된 색을 run에 되살림 (가능한 경우만)"""
    try:
        if snap and snap.get("type") == "rgb" and snap.get("rgb") is not None:
            run.font.color.rgb = snap["rgb"]
        elif snap and snap.get("type") == "theme" and snap.get("theme") is not None:
            run.font.color.theme_color = snap["theme"]
    except Exception:
        pass

# =========================
# CID용 계산/마커 포맷
# =========================
def load_cid(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    df = pd.read_csv(file_path) if ext==".csv" else pd.read_excel(file_path)
    if "Month" not in df.columns:
        raise KeyError("CID 파일에 Month 컬럼이 필요합니다.")
    df["Month"] = pd.to_datetime(df["Month"], errors="coerce")
    df["year"]  = df["Month"].dt.year
    df["month"] = df["Month"].dt.month
    return df

def safe_pct(a, b):
    if b == 0 or pd.isna(b): return np.nan
    return (a - b) / b * 100

def calc_metrics(sub):
    gms   = sub["GMS"].sum()
    gv    = sub["GV"].sum()
    units = sub["Units"].sum()
    ba    = sub["BA"].sum() if "BA" in sub.columns else 0
    asp   = gms / units if units > 0 else np.nan
    cr    = units / gv * 100 if gv > 0 else np.nan
    gcpb  = gms / ba if (ba is not None and ba > 0) else np.nan
    return {"gms": gms, "gv": gv, "units": units, "ba": ba, "asp": asp, "cr": cr, "gcpb": gcpb}

def fmt_gms(v):
    if pd.isna(v): return "N/A"
    if v >= 1_000_000: return f"{v/1_000_000:.1f}M"
    if v >= 1_000:     return f"{v/1_000:.0f}K"
    return f"{v:.0f}"

def fmt_k(v):
    if pd.isna(v) or v <= 0: return "N/A"
    return f"{v/1_000:.0f}K"

def fmt_pct(v, dec=1):
    if pd.isna(v): return "N/A"
    return f"{float(v):.{dec}f}%"

def style_percent_text_and_color(text, run):
    m = re.search(r'(-?\d+(?:\.\d+)?)\s*%', text or "")
    if not m:
        return text
    try:
        val = float(m.group(1))
    except ValueError:
        return text
    # 색상 지정은 try로 감싸 안전 처리
    if val > 0:
        try: run.font.color.rgb = RGBColor(0, 176, 80)
        except Exception: pass
        if "▲" not in text and "▼" not in text: return text + " ▲"
    elif val < 0:
        try: run.font.color.rgb = RGBColor(192, 0, 0)
        except Exception: pass
        if "▲" not in text and "▼" not in text: return text + " ▼"
    return text

def month_back(year, month, n):
    y, m = year, month
    for _ in range(n):
        if m == 1: y -= 1; m = 12
        else: m -= 1
    return y, m

def metric_key_from_marker(marker):
    m = marker.lower()
    if   "gms"   in m: return "gms"
    elif "gv"    in m: return "gv"
    elif "units" in m: return "units"
    elif "ba"    in m: return "ba"
    elif "asp"   in m: return "asp"
    elif "cr"    in m: return "cr"
    elif "gcpb"  in m: return "gcpb"
    elif "ipi"   in m: return "ipi"
    elif "excess" in m: return "excess"
    elif "avgwoc" in m or "avg.woc" in m or "avg_woc" in m: return "avgwoc"
    elif "fbagms" in m or "fba gms" in m: return "fbagms"
    elif "awas" in m: return "awas"
    elif "awagv" in m: return "awagv"
    return None

def _format_by_marker(key, metrics):
    if key == "gms":   return fmt_gms(metrics["gms"])
    if key == "gv":    return fmt_k(metrics["gv"])
    if key == "units": return fmt_k(metrics["units"])
    if key == "ba":    return f"{metrics['ba']:.0f}" if (not pd.isna(metrics["ba"]) and metrics["ba"]>0) else "N/A"
    if key == "asp":   return f"{metrics['asp']:.1f}" if not pd.isna(metrics["asp"]) else "N/A"
    if key == "cr":    return fmt_pct(metrics["cr"])
    if key == "gcpb":  return f"{metrics['gcpb']:.1f}" if not pd.isna(metrics["gcpb"]) else "N/A"
    return "N/A"

# ---------- (카테고리 텍스트 마커: 기존 유지) ----------
def _get_topcat_text_marker(marker):
    s = marker.lower().replace(" ", "")
    if "top1category" in s:      return TOPCAT_METRICS.get("top1_category", "N/A")
    if "top1gmsportion" in s:    return fmt_pct(TOPCAT_METRICS.get("top1_portion", np.nan))
    if "top1gmsgrowth" in s or "top1monthlygrowth" in s:     # 추가된 조건
        return fmt_pct(TOPCAT_METRICS.get("top1_growth", np.nan))
    if "top2category" in s:      return TOPCAT_METRICS.get("top2_category", "N/A")
    if "top2gmsportion" in s:    return fmt_pct(TOPCAT_METRICS.get("top2_portion", np.nan))
    if "top2gmsgrowth" in s or "top2monthlygrowth" in s:     # 추가된 조건
        return fmt_pct(TOPCAT_METRICS.get("top2_growth", np.nan))
    return None
# ---------- (ASIN 텍스트 마커: 추가) ----------
def _norm_colname(s: str) -> str:
    return str(s).lower().replace(" ", "").replace("/", "")

def _get_col_any(df, *names):
    """df에서 여러 이름 후보 중 존재하는 실제 컬럼명을 찾아 반환"""
    norm_map = {_norm_colname(c): c for c in df.columns}
    for nm in names:
        key = _norm_colname(nm)
        if key in norm_map:
            return norm_map[key]
    return None  # KeyError 대신 None 반환

def _fmt_or_na(v, dec=1):
    return "N/A" if pd.isna(v) else f"{float(v):.{dec}f}"

def compute_topasin_metrics(df_asin, df_cid):
    """
    최신월(anchor) 기준 Top1/Top2 Child ASIN을 찾고
    - ASIN 코드
    - 매출 비중(해당월 Top 전체 대비 %)
    - MoM 성장률(해당 ASIN GMS 기준)
    을 TOPASIN_METRICS에 저장.
    """
    global TOPASIN_METRICS
    TOPASIN_METRICS = {}

    if "Month" not in df_asin.columns:
        return
    da = df_asin.copy()
    da["Month"] = pd.to_datetime(da["Month"], errors="coerce")
    da = da.dropna(subset=["Month"])
    if da.empty:
        return

    asin_col = _get_col_any(da, "ASIN", "Child ASIN", "ChildASIN")
    gms_col  = _get_col_any(da, "GMS")

    # 앵커월: CID 최신월 우선, 없으면 ASIN 최신월
    try:
        anchor_dt = pd.to_datetime(df_cid["Month"], errors="coerce").max()
        if pd.isna(anchor_dt):
            anchor_dt = da["Month"].max()
    except Exception:
        anchor_dt = da["Month"].max()

    cur = da[da["Month"].dt.to_period("M") == anchor_dt.to_period("M")]
    if cur.empty:
        cur = da
        anchor_dt = da["Month"].max()

    total_gms = float(cur[gms_col].sum())
    prev_dt = (anchor_dt - pd.offsets.MonthBegin(1))

    top = (cur.groupby(asin_col, dropna=False)[gms_col]
             .sum()
             .sort_values(ascending=False))
    asins = [str(x) for x in top.index.tolist()[:2]]

    for rank, asin in enumerate(asins, start=1):
        cur_gms  = float(cur[cur[asin_col]==asin][gms_col].sum())
        prev_gms = float(da[(da[asin_col]==asin) &
                            (da["Month"].dt.to_period("M")==prev_dt.to_period("M"))][gms_col].sum())
        portion  = (cur_gms/total_gms*100) if total_gms>0 else np.nan
        growth   = safe_pct(cur_gms, prev_gms)

        TOPASIN_METRICS[f"top{rank}_asin"]    = asin
        TOPASIN_METRICS[f"top{rank}_portion"] = portion
        TOPASIN_METRICS[f"top{rank}_growth"]  = growth








def _get_topasin_text_marker(marker: str):
    """
    PPT 마커 치환:
      {Top1 ASIN} / {Top2 ASIN}
      {Top1 ASIN portion} / {Top2 ASIN portion}        -> 12.3   (숫자만)
      {Top1 ASIN growth}  / {Top2 ASIN growth}         -> 3.1    (절대값 숫자만)
    퍼센트 자동 포함 버전:
      {Top1 ASIN portion%} / {Top1 ASIN growth%}       -> "12.3%" / "-3.1%"
    """
    s = marker.lower().strip().replace(" ", "")
    for r in (1, 2):
        if s == f"top{r}asin":
            return TOPASIN_METRICS.get(f"top{r}_asin", "N/A")
        if s == f"top{r}asinportion":
            return _fmt_or_na(TOPASIN_METRICS.get(f"top{r}_portion", np.nan))
        if s == f"top{r}asingrowth":
            g = TOPASIN_METRICS.get(f"top{r}_growth", np.nan)
            return _fmt_or_na(abs(g))
        if s == f"top{r}asinportion%":
            return fmt_pct(TOPASIN_METRICS.get(f"top{r}_portion", np.nan))
        if s == f"top{r}asingrowth%":
            return fmt_pct(TOPASIN_METRICS.get(f"top{r}_growth", np.nan))
    return None
def extract_value(df, marker):
    s = (marker or "").strip().lower()
    s = s.replace('@', '_')
    # ===== 1순위: CID 월별 값 마커 (IPI, Excess, WoC, GMS, FBA GMS 등) =====
    mcid = re.match(r'^(ipi|excess|woc|gms|fbagms|fbagms_pct)_mm(?:-(\d+))?_(\d{4})$', s)
    if mcid:
        kind = mcid.group(1)                 
        back = int(mcid.group(2) or 0)       
        yreq = int(mcid.group(3))            

        if df[df["year"] == yreq].empty:
            return "N/A"
        latest_m = int(df[df["year"] == yreq]["month"].max())
        ty, tm = month_back(yreq, latest_m, back)

        sub = df[(df["year"] == ty) & (df["month"] == tm)]
        if sub.empty:
            return "N/A"

        def _col(*cands):
            return _get_col_any(df, *cands)

        try:
            if kind == "ipi":
                col = _col("IPI Score", "IPI")
                if not col: return "N/A"
                val = pd.to_numeric(sub[col], errors='coerce').mean()
                return "N/A" if pd.isna(val) else f"{val:.0f}"

            if kind == "excess":
                col = _col("Excess PCT", "Excess", "Excess%")
                if not col: return "N/A"
                val = pd.to_numeric(sub[col], errors='coerce').mean()
                return fmt_pct(val)

            if kind == "woc":
                col = _col("Avg WOC", "Avg. WoC", "Avg_WoC", "avgwoc", "WoC")
                if not col: return "N/A"
                val = pd.to_numeric(sub[col], errors='coerce').mean()
                return "N/A" if pd.isna(val) else f"{val:.1f}"

            if kind == "gms":
                col = _col("GMS")
                if not col: return "N/A"
                val = pd.to_numeric(sub[col], errors='coerce').sum()
                return fmt_k(val)

            if kind == "fbagms":
                col = _col("FBA GMS", "FBA_GMS", "FBA GMS Sales")
                if not col: return "N/A"
                val = pd.to_numeric(sub[col], errors='coerce').sum()
                return fmt_k(val)

            if kind == "fbagms_pct":
                col_f = _col("FBA GMS", "FBA_GMS", "FBA GMS Sales")
                col_g = _col("GMS")
                if not col_f or not col_g: return "N/A"
                f = pd.to_numeric(sub[col_f], errors='coerce').sum()
                g = pd.to_numeric(sub[col_g], errors='coerce').sum()
                pct = (f / g * 100) if g and g > 0 else np.nan
                return fmt_pct(pct)

        except Exception as e:
            print(f"CID 마커 처리 오류 {marker}: {e}")
            return "N/A"

    # ===== 2순위: 프로모션 마커 =====
    mpromo = re.match(r'^(promo|promo_pct)_mm(?:-(\d+))?_(\d{4})$', s)
    if mpromo:
        kind = mpromo.group(1)
        back = int(mpromo.group(2) or 0)
        yreq = int(mpromo.group(3))

        if df[df["year"] == yreq].empty:
            return "N/A"
        latest_m = int(df[df["year"] == yreq]["month"].max())
        ty, tm = month_back(yreq, latest_m, back)

        sub = df[(df["year"] == ty) & (df["month"] == tm)]
        if sub.empty:
            return "N/A"

        try:
            total_promo = 0
            
            promo_cols = [
                ("BD OPS", "BDOPS", "Best Deal"),
                ("LD OPS", "LDOPS", "Lightning Deal"),
                ("DOTD OPS", "DOTDOPS", "Deal of The Day"),
                ("Mario OPS", "MarioOPS", "Prime Exclusive"),
                ("Coupon OPS", "CouponOPS", "Coupon")
            ]
            
            for col_candidates in promo_cols:
                col = _get_col_any(df, *col_candidates)
                if col:
                    val = pd.to_numeric(sub[col], errors='coerce').sum()
                    total_promo += val if not pd.isna(val) else 0

            if kind == "promo":
                return fmt_k(total_promo)
            
            elif kind == "promo_pct":
                gms_col = _get_col_any(df, "GMS")
                if gms_col:
                    gms = pd.to_numeric(sub[gms_col], errors='coerce').sum()
                    pct = (total_promo / gms * 100) if gms > 0 else np.nan
                    return fmt_pct(pct)
                return "N/A"

        except Exception as e:
            print(f"프로모션 마커 처리 오류: {e}")
            return "N/A"
    
    # ===== 3순위: 날짜 라벨 전용 =====
    mdate = re.match(r'^mm(?:-(\d+))?(?:_(\d{4}))?$', s)
    if mdate:
        back = int(mdate.group(1) or 0)
        base_year = int(mdate.group(2)) if mdate.group(2) else int(df['year'].max())
        if df[df['year'] == base_year].empty:
            return "N/A"
        base_month = int(df[df['year'] == base_year]['month'].max())
        ty, tm = month_back(base_year, base_month, back)
        return f"{ty}-{tm:02d}"
    
    # ===== 4순위: Top ASIN 텍스트 마커 =====
    val_asin = _get_topasin_text_marker(marker)
    if val_asin is not None:
        return val_asin

    # ===== 5순위: Top Category 텍스트 마커 =====  
    val_top = _get_topcat_text_marker(marker)
    if val_top is not None:
        return val_top

    # ===== 이하 기존 로직 =====
    marker = s
    year_match = re.search(r"\d{4}", marker)
    year = int(year_match.group()) if year_match else int(df['year'].max())
    if df[df["year"] == year].empty:
        return "N/A"

    latest_month = int(df[df["year"] == year]["month"].max())
    anchor_year  = int(df["year"].max())
    anchor_month = int(df[df["year"] == anchor_year]["month"].max())

    mm_offset = None
    m = re.search(r"mm-(\d+)", marker)
    if m: mm_offset = int(m.group(1))

    def pick_month_subset(y, m_):
        return df[(df["year"] == y) & (df["month"] == m_)]

    try:
        if "y-1" in marker:
            prev_year = year - 1
            if not any(k in marker for k in ["gms","gv","units","ba","asp","cr","gcpb","ipi","excess","avgwoc","fbagms","awas","awagv"]):
                return f"{prev_year}-{str(latest_month).zfill(2)}"
            sub = pick_month_subset(prev_year, latest_month)
            if sub.empty: return "N/A"
            key = metric_key_from_marker(marker)
            if key in ("ipi","excess","avgwoc","fbagms","awas","awagv"): return "N/A"
            return _format_by_marker(key, calc_metrics(sub))

        if "m-1" in marker:
            prev_month = 12 if latest_month == 1 else latest_month - 1
            prev_year  = year - 1 if latest_month == 1 else year
            if not any(k in marker for k in ["gms","gv","units","ba","asp","cr","gcpb","ipi","excess","avgwoc","fbagms","awas","awagv"]):
                return f"{prev_year}-{str(prev_month).zfill(2)}"
            sub = pick_month_subset(prev_year, prev_month)
            if sub.empty: return "N/A"
            key = metric_key_from_marker(marker)
            if key in ("ipi","excess","avgwoc","fbagms","awas","awagv"): return "N/A"
            return _format_by_marker(key, calc_metrics(sub))

        if mm_offset is not None:
            ty, tm = month_back(year, latest_month, mm_offset)
            if not any(k in marker for k in ["gms","gv","units","ba","asp","cr","gcpb","ipi","excess","avgwoc","fbagms","awas","awagv"]):
                return f"{ty}-{str(tm).zfill(2)}"
            sub = pick_month_subset(ty, tm)
            if sub.empty: return "N/A"
            key = metric_key_from_marker(marker)
            if key in ("ipi","excess","avgwoc","fbagms","awas","awagv"): return "N/A"
            return _format_by_marker(key, calc_metrics(sub))

        if "mm" in marker and "mom" not in marker and "yoy" not in marker:
            if not any(k in marker for k in ["gms","gv","units","ba","asp","cr","gcpb","ipi","excess","avgwoc","fbagms","awas","awagv"]):
                return f"{year}-{str(latest_month).zfill(2)}"
            sub = pick_month_subset(year, latest_month)
            if sub.empty: return "N/A"
            key = metric_key_from_marker(marker)
            if key in ("ipi","excess","avgwoc","fbagms","awas","awagv"): return "N/A"
            return _format_by_marker(key, calc_metrics(sub))

        if "ytd" in marker and "yoy" not in marker:
            cur = calc_metrics(df[(df["year"] == year) & (df["month"] <= anchor_month)])
            key = metric_key_from_marker(marker)
            if key in ("ipi","excess","avgwoc","fbagms","awas","awagv"): return "N/A"
            return _format_by_marker(key, cur)

        if "ytd" in marker and "yoy" in marker:
            key = metric_key_from_marker(marker)
            if key in ("ipi","excess","avgwoc","fbagms","awas","awagv"): return "N/A"
            cur  = calc_metrics(df[(df["year"] == year)   & (df["month"] <= anchor_month)])
            prev = calc_metrics(df[(df["year"] == year-1) & (df["month"] <= anchor_month)])
            return fmt_pct(safe_pct(cur[key], prev[key]) if prev[key] > 0 else np.nan)

        if "mm" in marker and "yoy" in marker:
            key = metric_key_from_marker(marker)
            if key in ("ipi","excess","avgwoc","fbagms","awas","awagv"): return "N/A"
            cur  = calc_metrics(df[(df["year"] == year)   & (df["month"] == latest_month)])
            prev = calc_metrics(df[(df["year"] == year-1) & (df["month"] == latest_month)])
            return fmt_pct(safe_pct(cur[key], prev[key]) if prev[key] > 0 else np.nan)

        if "mm" in marker and "mom" in marker:
            key = metric_key_from_marker(marker)
            if key in ("ipi","excess","avgwoc","fbagms","awas","awagv"): return "N/A"
            cur = calc_metrics(df[(df["year"] == year) & (df["month"] == latest_month)])
            prev_month = 12 if latest_month == 1 else latest_month-1
            prev_year  = year - 1 if latest_month == 1 else year
            prev = calc_metrics(df[(df["year"] == prev_year) & (df["month"] == prev_month)])
            return fmt_pct(safe_pct(cur[key], prev[key]) if prev[key] > 0 else np.nan)

    except Exception as e:
        print(f"⚠️ 마커 처리 실패: {marker} → {e}")
        return "N/A"

    return "N/A"



# =========================
# PPT 텍스트/테이블 치환
# =========================
def iter_all_shapes(shapes):
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in iter_all_shapes(sh.shapes):
                yield sub
        else:
            yield sh

def replace_text_preserve_style(shape, df):
    if not shape.has_text_frame:
        return

    # ---- local helpers (이 함수 안에서만 사용) ----
    def _norm(s):  # normalize col name
        return str(s).lower().replace(" ", "").replace("/", "")

    def _find_col_local(d, *cands):
        mp = {_norm(c): c for c in d.columns}
        for c in cands:
            k = _norm(c)
            if k in mp:
                return mp[k]
        return None

    def _compute_conversion(year: int, offset: int = 0):
        """요청 연도 최신월에서 offset개월 과거로 이동한 달의 (AWAS/BA*100) 반환"""
        if df[df["year"] == year].empty:
            return np.nan

        ba_col   = _find_col_local(df, "BA", "Buyable ASIN", "BuyableASIN")
        awas_col = _find_col_local(df, "AWAS", "AW AS", "AS w/ Sales", "AW w/ Sales",
                                   "AS with Sales", "AWAS Count")
        if not ba_col or not awas_col:
            return np.nan

        latest_m = int(df.loc[df["year"] == year, "month"].max())
        ty, tm = month_back(year, latest_m, offset)

        sub = df[(df["year"] == ty) & (df["month"] == tm)]
        if sub.empty:
            return np.nan

        # "8개" 같은 텍스트도 처리
        ba   = sub[ba_col].map(lambda v: parse_number_any(v)).sum()
        awas = sub[awas_col].map(lambda v: parse_number_any(v)).sum()

        return (awas / ba * 100.0) if (pd.notna(ba) and ba > 0) else np.nan
    # ---------------------------------------------

    for paragraph in shape.text_frame.paragraphs:
        # 현재 문단의 전체 텍스트(겹치는 마커 선치환 위해 run 합침)
        full_text = "".join(run.text for run in paragraph.runs)

        # 1) {awas_mm_YYYY} / {ba_mm_YYYY} 또는 역순  → 해당 연도 최신월 기준 Conversion%
        for pat in (r"\{awas_mm_(\d{4})\}\s*/\s*\{ba_mm_\1\}",
                    r"\{ba_mm_(\d{4})\}\s*/\s*\{awas_mm_\1\}"):
            while True:
                m = re.search(pat, full_text, flags=re.I)
                if not m:
                    break
                year = int(m.group(1))
                conv_val = _compute_conversion(year, 0)
                full_text = full_text[:m.start()] + fmt_pct(conv_val) + full_text[m.end():]

        # 2) {awas_mm-N_YYYY} / {ba_mm-N_YYYY} 또는 역순  → 최신월에서 N개월 전 Conversion%
        for pat in (r"\{awas_mm-(\d+)_(\d{4})\}\s*/\s*\{ba_mm-\1_\2\}",
                    r"\{ba_mm-(\d+)_(\d{4})\}\s*/\s*\{awas_mm-\1_\2\}"):
            while True:
                m = re.search(pat, full_text, flags=re.I)
                if not m:
                    break
                off = int(m.group(1))
                year = int(m.group(2))
                conv_val = _compute_conversion(year, off)
                full_text = full_text[:m.start()] + fmt_pct(conv_val) + full_text[m.end():]

        # 3) 나머지 {marker}들은 일반 치환 로직으로 처리하기 위해 분해
        parts, last_idx = [], 0
        for match in re.finditer(r"\{([^}]+)\}", full_text):
            marker = match.group(1)
            parts.append((full_text[last_idx:match.start()], None))   # 일반 텍스트
            parts.append((match.group(0), marker))                    # 마커 텍스트
            last_idx = match.end()
        parts.append((full_text[last_idx:], None))

        # 4) 첫 run 스타일 스냅샷
        first_run = paragraph.runs[0] if paragraph.runs else None
        font_name = first_run.font.name if first_run else "Arial"
        font_size = first_run.font.size if (first_run and first_run.font.size) else Pt(12)
        bold      = first_run.font.bold if first_run else None
        italic    = first_run.font.italic if first_run else None
        underline = first_run.font.underline if first_run else None
        color_snap = _snapshot_color(first_run.font.color) if first_run else {"type": None}

        # 5) 문단 재작성(스타일 유지)
        paragraph.clear()
        for text, marker in parts:
            run = paragraph.add_run()
            run.font.name, run.font.size = font_name, font_size
            run.font.bold, run.font.italic, run.font.underline = bold, italic, underline
            _apply_color(run, color_snap)

            if marker:
                val = extract_value(df, marker)
                run.text = val
                run.text = style_percent_text_and_color(run.text, run)
            else:
                run.text = text



def process_ppt_markers(prs, df):
    for slide in prs.slides:
        for shape in iter_all_shapes(slide.shapes):
            replace_text_preserve_style(shape, df)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                matches = re.findall(r"\{([^}]+)\}", run.text)
                                for marker in matches:
                                    val = extract_value(df, marker)
                                    run.text = run.text.replace("{"+marker+"}", val if val else "N/A")
                                run.text = style_percent_text_and_color(run.text, run)

# =========================
# YTD 테이블 자동 채움
# =========================
def find_title_shape(slide, regex):
    pat = re.compile(regex)
    for sh in iter_all_shapes(slide.shapes):
        if getattr(sh, "has_text_frame", False) and sh.has_text_frame:
            txt = "\n".join(p.text for p in sh.text_frame.paragraphs)
            if pat.search(txt): 
                return sh
    return None

def center(shape):
    return (shape.left + shape.width/2, shape.top + shape.height/2)

def nearest_table(slide, anchor_shape):
    ax, ay = center(anchor_shape)
    best, best_d2 = None, None
    for sh in iter_all_shapes(slide.shapes):
        if hasattr(sh, "has_table") and sh.has_table:
            cx, cy = center(sh)
            d2 = (cx-ax)*(cx-ax) + (cy-ay)*(cy-ay)
            if best is None or d2 < best_d2:
                best, best_d2 = sh, d2
    return best.table if best else None

def write_cell(cell, text):
    tf = cell.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        r0 = tf.paragraphs[0].runs[0]
        font_name = r0.font.name
        font_size = r0.font.size or Pt(12)
        color_snap = _snapshot_color(r0.font.color)
    else:
        font_name, font_size = "Arial", Pt(12)
        color_snap = {"type": None}
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.font.name, run.font.size = font_name, font_size
    _apply_color(run, color_snap)
    run.text = text
    run.text = style_percent_text_and_color(run.text, run)

def fill_ytd_table(table, df):
    anchor_year  = int(df["year"].max())
    anchor_month = int(df[df["year"] == anchor_year]["month"].max())
    years = [anchor_year, anchor_year-1, anchor_year-2]

    def ytd(y):
        sub = df[(df["year"] == y) & (df["month"] <= anchor_month)]
        if sub.empty:
            return {"gms": np.nan, "gv": np.nan, "units": np.nan, "cr": np.nan}
        gms   = sub["GMS"].sum()
        gv    = sub["GV"].sum()
        units = sub["Units"].sum()
        cr    = (units / gv * 100) if gv > 0 else np.nan
        return {"gms": gms, "gv": gv, "units": units, "cr": cr}

    yvals = {y: ytd(y) for y in years}

    def yoy(cur, prev, key):
        a = yvals[cur][key]; b = yvals[prev][key]
        return safe_pct(a, b) if (not pd.isna(b) and b > 0) else np.nan

    for ridx, y in enumerate(years, start=1):
        if ridx >= len(table.rows): break
        try:
            write_cell(table.cell(ridx, 0), f"{y} YTD (≤ {str(anchor_month).zfill(2)})")
        except Exception:
            pass
        v = yvals[y]
        write_cell(table.cell(ridx, 1), fmt_gms(v["gms"]))
        write_cell(table.cell(ridx, 2), fmt_k(v["gv"]))
        write_cell(table.cell(ridx, 3), fmt_k(v["units"]))
        write_cell(table.cell(ridx, 4), fmt_pct(v["cr"]))
        if y-1 in yvals:
            write_cell(table.cell(ridx, 5), fmt_pct(yoy(y, y-1, "gms")))
            write_cell(table.cell(ridx, 6), fmt_pct(yoy(y, y-1, "gv")))
            write_cell(table.cell(ridx, 7), fmt_pct(yoy(y, y-1, "units")))
            write_cell(table.cell(ridx, 8), fmt_pct(yoy(y, y-1, "cr")))
        else:
            for c in (5,6,7,8):
                write_cell(table.cell(ridx, c), "N/A")
def build_top10_asin_table(df_asin, df_cid):
    df = df_asin.copy()
    df["Month"] = pd.to_datetime(df["Month"], errors="coerce")
    df = df.dropna(subset=["Month"])
    asin_col  = _get_col_any(df, "ASIN", "Child ASIN", "ChildASIN")
    gms_col   = _get_col_any(df, "GMS")
    gv_col    = _get_col_any(df, "GV")
    units_col = _get_col_any(df, "Units")

    # 앵커월
    anchor_dt = pd.to_datetime(df_cid["Month"], errors="coerce").max()
    cur = df[df["Month"].dt.to_period("M")==anchor_dt.to_period("M")]
    if cur.empty:
        cur = df
        anchor_dt = df["Month"].max()

    total_gms = cur[gms_col].sum()

    # Top10 추출
    top10 = (cur.groupby(asin_col)[gms_col]
                .sum()
                .sort_values(ascending=False)
                .head(10)
                .reset_index())

    out = []
    for _, row in top10.iterrows():
        asin = row[asin_col]
        gms  = row[gms_col]
        portion = gms/total_gms*100 if total_gms>0 else np.nan

        # 전체 데이터 기준 시계열 추출
        sub = df[df[asin_col]==asin].copy()
        sub["year"] = sub["Month"].dt.year
        sub["month"] = sub["Month"].dt.month
        sub["ASP"] = np.where(sub[units_col]>0, sub[gms_col]/sub[units_col], np.nan)
        sub["CR"]  = np.where(sub[gv_col]>0, sub[units_col]/sub[gv_col]*100, np.nan)
        agg = finalize_year_month(sub.groupby(["year","month"])
                                   .agg(GMS=(gms_col,"sum"),
                                        GV=(gv_col,"sum"),
                                        Units=(units_col,"sum"),
                                        ASP=("ASP","mean"),
                                        CR=("CR","mean"))
                                   .reset_index())

        # 최신월과 이전월 비교 (MoM)
        cur_row = agg[agg["year"]==anchor_dt.year]
        cur_row = cur_row[cur_row["month"]==anchor_dt.month]
        prev_row = agg[(agg["year"]==month_back(anchor_dt.year, anchor_dt.month, 1)[0]) &
                       (agg["month"]==month_back(anchor_dt.year, anchor_dt.month, 1)[1])]

        def pct(cur, prev): return safe_pct(cur, prev) if prev>0 else np.nan

        gms_mom   = pct(cur_row["GMS"].sum(),   prev_row["GMS"].sum()) if not prev_row.empty else np.nan
        gv_mom    = pct(cur_row["GV"].sum(),    prev_row["GV"].sum()) if not prev_row.empty else np.nan
        units_mom = pct(cur_row["Units"].sum(), prev_row["Units"].sum()) if not prev_row.empty else np.nan
        asp_mom   = pct(cur_row["ASP"].mean(),  prev_row["ASP"].mean()) if not prev_row.empty else np.nan
        cr_mom    = pct(cur_row["CR"].mean(),   prev_row["CR"].mean()) if not prev_row.empty else np.nan

        # YoY
        prev_y = anchor_dt.year-1
        prev_row_y = agg[(agg["year"]==prev_y) & (agg["month"]==anchor_dt.month)]
        gms_yoy   = pct(cur_row["GMS"].sum(),   prev_row_y["GMS"].sum()) if not prev_row_y.empty else np.nan
        gv_yoy    = pct(cur_row["GV"].sum(),    prev_row_y["GV"].sum()) if not prev_row_y.empty else np.nan
        units_yoy = pct(cur_row["Units"].sum(), prev_row_y["Units"].sum()) if not prev_row_y.empty else np.nan
        asp_yoy   = pct(cur_row["ASP"].mean(),  prev_row_y["ASP"].mean()) if not prev_row_y.empty else np.nan
        cr_yoy    = pct(cur_row["CR"].mean(),   prev_row_y["CR"].mean()) if not prev_row_y.empty else np.nan

        out.append([asin, gms, portion, 
                    cur_row["GV"].sum(), cur_row["Units"].sum(),
                    cur_row["ASP"].mean(), cur_row["CR"].mean(),
                    gms_mom, gv_mom, units_mom, asp_mom, cr_mom,
                    gms_yoy, gv_yoy, units_yoy, asp_yoy, cr_yoy])
    return pd.DataFrame(out, columns=["ASIN","GMS","GMS%","GV","Units","ASP","CR",
                                      "GMS MoM","GV MoM","Units MoM","ASP MoM","CR MoM",
                                      "GMS YoY","GV YoY","Units YoY","ASP YoY","CR YoY"])

def fill_top10_asin_table(prs, df_asin, df_cid):
    table_df = build_top10_asin_table(df_asin, df_cid)
    for slide in prs.slides:
        title = find_title_shape(slide, r"Top10\s*Child\s*ASIN")
        if not title: continue
        table = nearest_table(slide, title)
        if not table: continue
        for r, (_, row) in enumerate(table_df.iterrows(), start=1):  # 첫 행은 헤더니까 row=1부터
            for c, val in enumerate(row):
                write_cell(table.cell(r, c), fmt_pct(val) if isinstance(val,float) and "pct" in table_df.columns[c].lower() else str(round(val,2) if isinstance(val,float) else val))


def fill_ytd_table_on_slides(prs, df):
    count = 0
    for slide in prs.slides:
        title = find_title_shape(slide, r"연간\s*누적\s*매출\s*트렌드")
        if not title: continue
        table = nearest_table(slide, title)
        if not table: continue
        try:
            fill_ytd_table(table, df); count += 1
        except Exception as e:
            print(f"YTD 테이블 채우기 실패: {e}")
    return count


# =========================
# 그래프 생성 (원본 로직 준수)
# =========================
def create_line_by_year(dates, values, title, y_label, graph_name,
                        unit="none", decimal=1, percentage=False, annotate_year=None):
    df = monthly_agg(dates, values)
    if df.empty:
        return None
    df["value_conv"] = df["value"]
    _set_korean_font_if_possible()
    fig = plt.figure(figsize=(10, 6)); ax = plt.gca(); _bi_theme(ax)
    colors_cycle = [PALETTE["primary"], PALETTE["green"], PALETTE["orange"], PALETTE["purple"], PALETTE["red"]]
    for idx, y in enumerate(sorted(df["year"].unique())):
        sub = df[df["year"] == y].sort_values("month")
        xs = sub["month"].to_numpy(); ys = sub["value_conv"].to_numpy()
        ax.plot(xs, ys, linewidth=2.5, marker="o", markersize=5, color=colors_cycle[idx%len(colors_cycle)], label=f"{y}년")
        if annotate_year and y == annotate_year:
            txt = f"{ys[-1]:.{decimal}f}{'%' if percentage else ''}"
            _label_last(ax, xs, ys, txt)
    ax.set_title(title, fontsize=16, color=PALETTE["dark"])
    ax.set_xlabel("월", fontsize=11, color=PALETTE["dark"])
    ax.set_ylabel(y_label, fontsize=11, color=PALETTE["dark"])
    ax.set_xticks(range(1,13)); ax.set_xticklabels(MONTH_LABELS)
    if unit=="K": ax.yaxis.set_major_formatter(_yfmt_k(decimal))
    elif percentage: ax.yaxis.set_major_formatter(_yfmt_decimal(decimal, "%"))
    else: ax.yaxis.set_major_formatter(_yfmt_decimal(decimal))
    ax.legend(loc="upper left", fontsize=10, frameon=False)
    fig.tight_layout()
    return _save_fig(fig, graph_name)

def create_combo_ba_awagv_awas(dates, ba_values, awagv_values, awas_values, graph_name, title):
    # 월별 BA 막대 + Discoverability(= AWAGV/BA %) / Conversion(= AWAS/BA %) 라인
    rows=[]; n=len(dates)
    for i in range(n):
        dt = parse_date_any(dates[i])
        if dt is None: 
            continue
        ba   = parse_number_any(ba_values[i])    if i < len(ba_values)    else None
        awag = parse_number_any(awagv_values[i]) if i < len(awagv_values) else None
        awas = parse_number_any(awas_values[i])  if i < len(awas_values)  else None
        if ba is None or ba <= 0:
            continue
        disc = (awag/ba*100) if awag is not None else np.nan   # Discoverability %
        conv = (awas/ba*100) if awas is not None else np.nan   # Conversion %
        rows.append({"year": dt.year, "month": dt.month, "ba": ba, "disc": disc, "conv": conv})

    if not rows:
        return None

    df = pd.DataFrame(rows).groupby(["year","month"]).mean(numeric_only=True).reset_index()
    df = finalize_year_month(df, "year", "month")

    _set_korean_font_if_possible()
    fig = plt.figure(figsize=(10, 6)); ax1 = plt.gca(); _bi_theme(ax1)

    x = np.arange(len(df))

    # BA 막대
    bars = ax1.bar(x, df["ba"], alpha=0.9, color=PALETTE["ba_fill"], label="BA")
    ax1.set_ylabel("취급 상품 개수 (개)", color=PALETTE["dark"])

    # 보조축 라인들
    ax2 = ax1.twinx(); _bi_theme(ax2); ax2.set_yticks([])
    line_disc, = ax2.plot(x, df["disc"], "o-", linewidth=2.5, color=PALETTE["primary"],
                          label="Discoverability (AWAGV/BA%)")
    line_conv, = ax2.plot(x, df["conv"], "s-", linewidth=2.5, color=PALETTE["orange"],
                          label="Conversion (AWAS/BA%)")

    # ---------- Conversion 데이터 라벨: 2025년만 ----------
    target_year = 2025
    if target_year not in set(df["year"].astype(int)):
        # 만약 2025 데이터가 없다면 최신 연도만 표시(겹침 방지 목적 유지)
        target_year = int(df["year"].max())

    y_conv = df["conv"].to_numpy()
    for i in range(len(df)):
        if int(df.loc[i, "year"]) != target_year or pd.isna(y_conv[i]):
            continue
        # 주변 포인트와 겹치지 않게 위/아래로 번갈아 배치
        dy = 8 if (i % 2 == 0) else -10
        va = "bottom" if dy > 0 else "top"
        ax2.annotate(f"{y_conv[i]:.1f}%", (x[i], y_conv[i]),
                     textcoords="offset points", xytext=(0, dy),
                     ha="center", va=va, fontsize=8)

    # 제목/축/범례
    ax1.set_title(title, fontsize=16, color=PALETTE["dark"])
    ax1.set_xlabel("연월", color=PALETTE["dark"])
    ax1.set_xticks(x); ax1.set_xticklabels(df["date_str"], rotation=45, ha="right")
    ax1.legend([bars, line_disc, line_conv],
               ["BA", "Discoverability (AWAGV/BA%)", "Conversion (AWAS/BA%)"],
               loc="upper left", fontsize=9, frameon=False)

    fig.tight_layout()
    return _save_fig(fig, graph_name)



def create_ipi_combo_graph(dates, ipi_values, excess_pct_values, graph_name, title):
    rows = []
    for i in range(len(dates)):
        dt = parse_date_any(dates[i])
        ipi = parse_number_any(ipi_values[i]) if i<len(ipi_values) else None
        exc = parse_number_any(excess_pct_values[i], pct_to_100=True) if i<len(excess_pct_values) else None
        if dt is None or ipi is None: continue
        rows.append({"year":dt.year,"month":dt.month,"ipi":ipi,"excess":exc or np.nan})
    if not rows: return None
    df = pd.DataFrame(rows).groupby(["year","month"]).mean().reset_index()
    df = finalize_year_month(df, "year", "month")
    _set_korean_font_if_possible()
    fig = plt.figure(figsize=(10, 6)); ax1 = plt.gca(); _bi_theme(ax1)
    x = np.arange(len(df))
    bars = ax1.bar(x, df["ipi"], alpha=0.9, color=PALETTE["ba_fill"], label="IPI Score")
    ax1.set_ylabel("IPI Score", color=PALETTE["dark"])
    ax2 = ax1.twinx(); _bi_theme(ax2); ax2.set_yticks([])
    y2 = df["excess"].to_numpy()
    ax2.plot(x, y2, "o-", linewidth=2.5, color=PALETTE["orange"], label="Excess PCT")
    _label_last(ax2, x, y2, f"{y2[-1]:.1f}%")
    ax2.set_ylabel("Excess PCT (%)", color=PALETTE["dark"])
    ax1.set_title(title, fontsize=16, color=PALETTE["dark"])
    ax1.set_xlabel("연월", color=PALETTE["dark"])
    ax1.set_xticks(x); ax1.set_xticklabels(df["date_str"], rotation=45, ha="right")
    ax1.legend([bars], ["IPI Score"], loc="upper left", fontsize=10, frameon=False)
    fig.tight_layout()
    return _save_fig(fig, graph_name)

def create_merchandising_graph(dates, total_sales, bd_ops, ld_ops, dotd_ops, mario_ops, coupon_ops, graph_name, title):
    rows = []; n = len(dates)
    for i in range(n):
        dt = parse_date_any(dates[i])
        ts = parse_number_any(total_sales[i]) if i<len(total_sales) else None
        if dt is None or ts is None: continue
        bd = parse_number_any(bd_ops[i])    if i<len(bd_ops)    else 0
        ld = parse_number_any(ld_ops[i])    if i<len(ld_ops)    else 0
        do = parse_number_any(dotd_ops[i])  if i<len(dotd_ops)  else 0
        ma = parse_number_any(mario_ops[i]) if i<len(mario_ops) else 0
        cp = parse_number_any(coupon_ops[i])if i<len(coupon_ops)else 0
        total = (bd or 0)+(ld or 0)+(do or 0)+(ma or 0)+(cp or 0)
        pct = total/ts*100 if ts>0 else 0
        rows.append({"year":dt.year,"month":dt.month,"bd":bd or 0,"ld":ld or 0,"do":do or 0,"ma":ma or 0,"cp":cp or 0,"pct":pct})
    if not rows: return None
    df = pd.DataFrame(rows).groupby(["year","month"]).mean().reset_index()
    df = finalize_year_month(df, "year", "month")
    _set_korean_font_if_possible()
    fig = plt.figure(figsize=(12, 6)); ax1 = plt.gca(); _bi_theme(ax1)
    x = np.arange(len(df))
    bd, ld, do, ma, cp = df["bd"], df["ld"], df["do"], df["ma"], df["cp"]
    b1 = ax1.bar(x, bd, label="Best Deal", color=PALETTE["sp_fill"])
    b2 = ax1.bar(x, ld, bottom=bd, label="Lightning Deal", color=PALETTE["sb_fill"])
    b3 = ax1.bar(x, do, bottom=bd+ld, label="Deal of The Day", color=PALETTE["sd_fill"])
    b4 = ax1.bar(x, ma, bottom=bd+ld+do, label="Prime Exclusive Discount", color="#C7EBD0")
    b5 = ax1.bar(x, cp, bottom=bd+ld+do+ma, label="Coupon", color="#F8E3A2")
    ax1.set_ylabel("Merchandising OPS", color=PALETTE["dark"])
    ax2 = ax1.twinx(); _bi_theme(ax2); ax2.set_yticks([])
    y2 = df["pct"].to_numpy()
    ax2.plot(x, y2, "o-", linewidth=2.5, color=PALETTE["primary"], label="Merchandising OPS%")
    _label_last(ax2, x, y2, f"{y2[-1]:.1f}%")
    ax2.set_ylabel("Merchandising OPS%", color=PALETTE["dark"])
    ax1.set_title(title, fontsize=16, color=PALETTE["dark"])
    ax1.set_xlabel("연월", color=PALETTE["dark"])
    ax1.set_xticks(x); ax1.set_xticklabels(df["date_str"], rotation=45, ha="right")
    ax1.legend([b1,b2,b3,b4,b5], ["Best Deal","Lightning Deal","Deal of The Day","Prime Exclusive Discount","Coupon"],
               loc="upper left", fontsize=9, frameon=False)
    fig.tight_layout()
    return _save_fig(fig, graph_name)

def create_ads_tacos_graph(dates, total_sales, sp_spend, sb_spend, sd_spend, graph_name, title):
    rows=[]
    for i in range(len(dates)):
        dt = parse_date_any(dates[i])
        ts = parse_number_any(total_sales[i]) if i < len(total_sales) else None
        sp = parse_number_any(sp_spend[i])   if i < len(sp_spend)   else 0
        sb = parse_number_any(sb_spend[i])   if i < len(sb_spend)   else 0
        sd = parse_number_any(sd_spend[i])   if i < len(sd_spend)   else 0
        if dt is None or ts is None: continue
        total_ads = (sp or 0)+(sb or 0)+(sd or 0)
        tacos = total_ads/ts*100 if ts>0 else 0
        rows.append({"year":dt.year,"month":dt.month,"sp":sp or 0,"sb":sb or 0,"sd":sd or 0,"tacos":tacos})
    if not rows: return None
    df = pd.DataFrame(rows).groupby(["year","month"]).mean().reset_index()
    df = finalize_year_month(df, "year", "month")
    _set_korean_font_if_possible()
    fig = plt.figure(figsize=(12, 6)); ax1 = plt.gca(); _bi_theme(ax1)
    x = np.arange(len(df))
    b1 = ax1.bar(x, df["sp"], label="SP Spend", color=PALETTE["sp_fill"])
    b2 = ax1.bar(x, df["sb"], bottom=df["sp"], label="SB Spend", color=PALETTE["sb_fill"])
    b3 = ax1.bar(x, df["sd"], bottom=df["sp"]+df["sb"], label="SD Spend", color=PALETTE["sd_fill"])
    ax1.set_ylabel("Ads Spend", color=PALETTE["dark"])
    ax2 = ax1.twinx(); _bi_theme(ax2); ax2.set_yticks([])
    line_y = df["tacos"].to_numpy()
    ax2.plot(x, line_y, "o-", color=PALETTE["primary"], linewidth=2.5)
    _label_last(ax2, x, line_y, f"{line_y[-1]:.1f}%")
    ax1.set_title(title, fontsize=16, color=PALETTE["dark"])
    ax1.set_xlabel("연월", color=PALETTE["dark"]); ax1.set_xticks(x); ax1.set_xticklabels(df["date_str"], rotation=45, ha="right")
    ax1.legend([b1,b2,b3], ["SP Spend","SB Spend","SD Spend"], loc="upper left", fontsize=9, frameon=False)
    fig.tight_layout()
    return _save_fig(fig, graph_name)

def create_ads_impr_clicks_graph(dates, sp_imp, sp_clk, sb_imp, sb_clk, sd_imp, sd_clk, graph_name, title):
    rows=[]
    n=len(dates)
    for i in range(n):
        dt = parse_date_any(dates[i])
        if dt is None: continue
        sp_i = parse_number_any(sp_imp[i]) if i < len(sp_imp) else 0
        sp_c = parse_number_any(sp_clk[i]) if i < len(sp_clk) else 0
        sb_i = parse_number_any(sb_imp[i]) if i < len(sb_imp) else 0
        sb_c = parse_number_any(sb_clk[i]) if i < len(sb_clk) else 0
        sd_i = parse_number_any(sd_imp[i]) if i < len(sd_imp) else 0
        sd_c = parse_number_any(sd_clk[i]) if i < len(sd_clk) else 0
        rows.append({"year":dt.year,"month":dt.month,"sp_i":sp_i,"sb_i":sb_i,"sd_i":sd_i,
                     "clicks": (sp_c or 0)+(sb_c or 0)+(sd_c or 0)})
    if not rows: return None
    df = pd.DataFrame(rows).groupby(["year","month"]).mean().reset_index()
    df = finalize_year_month(df, "year", "month")
    _set_korean_font_if_possible()
    fig = plt.figure(figsize=(12, 6)); ax1 = plt.gca(); _bi_theme(ax1)
    x = np.arange(len(df))
    b1 = ax1.bar(x, df["sp_i"], label="SP Impression", color=PALETTE["sp_fill"])
    b2 = ax1.bar(x, df["sb_i"], bottom=df["sp_i"], label="SB Impression", color=PALETTE["sb_fill"])
    b3 = ax1.bar(x, df["sd_i"], bottom=df["sp_i"]+df["sb_i"], label="SD Impression", color=PALETTE["sd_fill"])
    ax1.set_ylabel("Impressions", color=PALETTE["dark"])
    ax2 = ax1.twinx(); _bi_theme(ax2); ax2.set_yticks([])
    clicks = df["clicks"].to_numpy()
    ax2.plot(x, clicks, "o-", color=PALETTE["primary"], linewidth=2.5, label="AD Clicks")
    _label_last(ax2, x, clicks, f"{clicks[-1]:.1f}")
    ax1.set_title(title, fontsize=16, color=PALETTE["dark"])
    ax1.set_xlabel("연월", color=PALETTE["dark"]); ax1.set_xticks(x); ax1.set_xticklabels(df["date_str"], rotation=45, ha="right")
    ax1.legend([b1,b2,b3], ["SP Impression","SB Impression","SD Impression"], loc="upper left", fontsize=9, frameon=False)
    fig.tight_layout()
    return _save_fig(fig, graph_name)

def create_three_line_pct(dates, a, b, c, labels, ylabel, graph_name, title):
    rows=[]
    n=len(dates)
    for i in range(n):
        dt = parse_date_any(dates[i])
        if dt is None: continue
        va = parse_number_any(a[i], pct_to_100=True) if i<len(a) else None
        vb = parse_number_any(b[i], pct_to_100=True) if i<len(b) else None
        vc = parse_number_any(c[i], pct_to_100=True) if i<len(c) else None
        rows.append({"year":dt.year,"month":dt.month, "A":va or 0, "B":vb or 0, "C":vc or 0})
    if not rows: return None
    df = pd.DataFrame(rows).groupby(["year","month"]).mean().reset_index()
    df = finalize_year_month(df, "year", "month")
    _set_korean_font_if_possible()
    fig = plt.figure(figsize=(12, 6)); ax = plt.gca(); _bi_theme(ax)
    x = np.arange(len(df))
    ax.plot(x, df["A"], "o-", linewidth=2.5, color=PALETTE["primary"], label=labels[0])
    ax.plot(x, df["B"], "s-", linewidth=2.5, color=PALETTE["orange"],  label=labels[1])
    ax.plot(x, df["C"], "^-", linewidth=2.5, color=PALETTE["purple"],  label=labels[2])
    _label_last(ax, x, df["A"].to_numpy(), f"{df['A'].iloc[-1]:.2f}%")
    ax.set_title(title, fontsize=16, color=PALETTE["dark"])
    ax.set_ylabel(ylabel, color=PALETTE["dark"]); ax.set_xlabel("연월", color=PALETTE["dark"])
    ax.set_xticks(x); ax.set_xticklabels(df["date_str"], rotation=45, ha="right")
    ax.yaxis.set_major_formatter(_yfmt_decimal(2, "%"))
    ax.legend(loc="upper left", fontsize=10, frameon=False)
    fig.tight_layout()
    return _save_fig(fig, graph_name)

def create_three_line_pct_nolabel(dates, a, b, c, labels, ylabel, graph_name, title):
    rows=[]; n=len(dates)
    for i in range(n):
        dt = parse_date_any(dates[i])
        if dt is None: continue
        va = parse_number_any(a[i], pct_to_100=True) if i<len(a) else None
        vb = parse_number_any(b[i], pct_to_100=True) if i<len(b) else None
        vc = parse_number_any(c[i], pct_to_100=True) if i<len(c) else None
        rows.append({'year':dt.year,'month':dt.month, 'A':va or 0, 'B':vb or 0, 'C':vc or 0})
    if not rows: return None
    df = pd.DataFrame(rows).groupby(['year','month']).mean().reset_index()
    df = finalize_year_month(df, 'year', 'month')
    _set_korean_font_if_possible()
    fig = plt.figure(figsize=(12,6)); ax = plt.gca(); _bi_theme(ax)
    x = np.arange(len(df))
    ax.plot(x, df['A'], 'o-', linewidth=2, label=labels[0])
    ax.plot(x, df['B'], 's-', linewidth=2, label=labels[1])
    ax.plot(x, df['C'], '^-', linewidth=2, label=labels[2])
    ax.set_title(title, fontsize=16)
    ax.set_ylabel(ylabel); ax.set_xlabel('연월')
    ax.set_xticks(x); ax.set_xticklabels(df['date_str'], rotation=45, ha='right')
    ax.legend(loc='upper left', fontsize=10, frameon=False)
    fig.tight_layout()
    return _save_fig(fig, graph_name)

def create_ads_sales_combo(dates, total_sales, sp_sales, sb_sales, sd_sales, graph_name, title):
    rows = []
    n = len(dates)
    for i in range(n):
        dt = parse_date_any(dates[i])
        if dt is None: 
            continue
        ts = parse_number_any(total_sales[i]) if i < len(total_sales) else None
        sp = parse_number_any(sp_sales[i])    if i < len(sp_sales)    else 0
        sb = parse_number_any(sb_sales[i])    if i < len(sb_sales)    else 0
        sd = parse_number_any(sd_sales[i])    if i < len(sd_sales)    else 0
        if ts is None:
            continue
        rows.append({
            "year": dt.year, "month": dt.month,
            "ts": ts or 0, "sp": sp or 0, "sb": sb or 0, "sd": sd or 0
        })

    if not rows:
        return None

    df = pd.DataFrame(rows).groupby(["year","month"]).sum(numeric_only=True).reset_index()
    df = finalize_year_month(df, "year", "month")
    df["ad_sales"] = df["sp"] + df["sb"] + df["sd"]
    df["ad_sales_pct"] = np.where(df["ts"] > 0, df["ad_sales"] / df["ts"] * 100, np.nan)

    _set_korean_font_if_possible()
    fig = plt.figure(figsize=(12, 6)); ax1 = plt.gca(); _bi_theme(ax1)

    x = np.arange(len(df))
    b1 = ax1.bar(x, df["sp"], label="SP Sales", color=PALETTE["sp_fill"])
    b2 = ax1.bar(x, df["sb"], bottom=df["sp"], label="SB Sales", color=PALETTE["sb_fill"])
    b3 = ax1.bar(x, df["sd"], bottom=df["sp"]+df["sb"], label="SD Sales", color=PALETTE["sd_fill"])

    ax1.set_ylabel("Ads Sales", color=PALETTE["dark"])

    ax2 = ax1.twinx(); _bi_theme(ax2); ax2.set_yticks([])
    y2 = df["ad_sales_pct"].to_numpy()
    ax2.plot(x, y2, "o-", linewidth=2.5, color=PALETTE["primary"], label="Ad sales%")
    if len(y2) > 0 and pd.notna(y2[-1]):
        _label_last(ax2, x, y2, f"{y2[-1]:.1f}%")

    ax1.set_title(title, fontsize=16, color=PALETTE["dark"])
    ax1.set_xlabel("연월", color=PALETTE["dark"])
    ax1.set_xticks(x); ax1.set_xticklabels(df["date_str"], rotation=45, ha="right")

    ax1.legend([b1,b2,b3], ["SP Sales","SB Sales","SD Sales"], loc="upper left", fontsize=9, frameon=False)
    fig.tight_layout()
    return _save_fig(fig, graph_name)

# =========================
# 그래프 11/12 (ASIN+CID)
# =========================
def _get_col(df, name):
    key = name.lower().replace(" ", "").replace("/", "")
    for c in df.columns:
        cc = str(c).lower().replace(" ", "").replace("/", "")
        if cc == key: return c
    raise KeyError(f"'{name}' 컬럼을 찾을 수 없습니다.")

def _fmt_big(v):
    if pd.isna(v): return "0"
    v = float(v)
    if v >= 1_000_000: return f"{v/1_000_000:.1f}M"
    if v >= 1_000:     return f"{v/1_000:.0f}K"
    return f"{v:.0f}"

def _heat_color(v, vmin=-100, vmax=100):
    import matplotlib.cm as cm, matplotlib.colors as colors
    if pd.isna(v): return "white"
    x = max(vmin, min(vmax, float(v)))
    norm = colors.Normalize(vmin=vmin, vmax=vmax)
    r,g,b,_ = cm.RdYlGn(norm(x))
    return (r, g, b)

def create_graph11_itkbn_dashboard(df_asin, df_monthly, sel_month=None, graph_name="Graph 11"):
    global TOPCAT_METRICS
    TOPCAT_METRICS = {}

    mcol = "Month" if "Month" in df_monthly.columns else None
    if mcol:
        df_monthly = df_monthly.copy()
        df_monthly["Month"] = pd.to_datetime(df_monthly["Month"], errors="coerce")
        anchor_dt = df_monthly["Month"].max()
        if sel_month:
            try: anchor_dt = pd.to_datetime(sel_month)
            except: pass
        month_str = anchor_dt.strftime("%Y-%m")
    else:
        month_str = "N/A"
        anchor_dt = pd.Timestamp.today().normalize()

    if "Month" not in df_asin.columns:
        raise KeyError("ASIN CSV에 Month 컬럼이 필요합니다. (YYYY-MM-DD 또는 YYYY-MM)")
    df_a = df_asin.copy()
    df_a["Month"] = pd.to_datetime(df_a["Month"], errors="coerce")
    df_a = df_a.dropna(subset=["Month"])

    col_itkbn = _get_col(df_a, "ITK/BN")
    col_gms   = _get_col(df_a, "GMS")
    col_gv    = _get_col(df_a, "GV")
    col_units = _get_col(df_a, "Units")

    cur = df_a[df_a["Month"].dt.to_period("M") == anchor_dt.to_period("M")]
    if cur.empty: cur = df_a
    donut_series = cur.groupby(col_itkbn, dropna=False)[col_gms].sum().sort_values(ascending=False)
    total_gms = float(donut_series.sum())

    top_list = donut_series.reset_index()
    top_list.columns = ["itkbn", "gms"]
    if not top_list.empty:
        top_list["portion"] = np.where(total_gms>0, top_list["gms"]/total_gms*100, np.nan)
        prev_dt = (anchor_dt - pd.offsets.MonthBegin(1))
        if len(top_list) >= 1:
            t1 = str(top_list.loc[0,"itkbn"])
            cur1 = cur[cur[col_itkbn]==t1][col_gms].sum()
            prev1 = df_a[(df_a[col_itkbn]==t1) & (df_a["Month"].dt.to_period("M")==prev_dt.to_period("M"))][col_gms].sum()
            g1 = safe_pct(cur1, prev1)
            TOPCAT_METRICS["top1_category"] = t1
            TOPCAT_METRICS["top1_portion"]  = float(top_list.loc[0,"portion"])
            TOPCAT_METRICS["top1_growth"]   = g1
        if len(top_list) >= 2:
            t2 = str(top_list.loc[1,"itkbn"])
            cur2 = cur[cur[col_itkbn]==t2][col_gms].sum()
            prev2 = df_a[(df_a[col_itkbn]==t2) & (df_a["Month"].dt.to_period("M")==prev_dt.to_period("M"))][col_gms].sum()
            g2 = safe_pct(cur2, prev2)
            TOPCAT_METRICS["top2_category"] = t2
            TOPCAT_METRICS["top2_portion"]  = float(top_list.loc[1,"portion"])
            TOPCAT_METRICS["top2_growth"]   = g2

    wc_series = cur.groupby(col_itkbn, dropna=False)[col_gms].sum().sort_values(ascending=False)
    wc_series = wc_series[wc_series>0].head(25)
    wc_labels = [str(k) if str(k)!="nan" else "Unknown" for k in wc_series.index]
    wc_vals   = wc_series.values.astype(float)

    TOP_N = 8
    donut_plot = donut_series.copy()
    if len(donut_plot) > TOP_N:
        others = donut_plot.iloc[TOP_N:].sum()
        donut_plot = donut_plot.iloc[:TOP_N]
        donut_plot.loc["Others"] = others
    donut_labels = [f"{k} ({_fmt_big(v)})" for k, v in zip(donut_plot.index, donut_plot.values)]
    donut_sizes  = donut_plot.values

    need_cols = {}
    for name in ["GMS","GV","Units"]:
        need_cols[name] = _get_col(df_monthly, name)
    ba_col = None
    for guess in ["BA","Buyable ASIN","BuyableASIN"]:
        try:
            ba_col = _get_col(df_monthly, guess); break
        except: pass

    mdf = df_monthly.copy()
    mdf["Month"] = pd.to_datetime(mdf["Month"], errors="coerce")
    mdf = mdf.dropna(subset=["Month"]).sort_values("Month")
    mdf["year"]  = mdf["Month"].dt.year
    mdf["month"] = mdf["Month"].dt.month

    agg = mdf.groupby(["year","month"]).agg({
        need_cols["GMS"]:"sum", need_cols["GV"]:"sum", need_cols["Units"]:"sum",
        ba_col if ba_col else need_cols["Units"]: "sum"
    }).reset_index()
    agg = agg.rename(columns={
        need_cols["GMS"]:"GMS", need_cols["GV"]:"GV", need_cols["Units"]:"Units",
        (ba_col if ba_col else need_cols["Units"]):"BA"
    })
    agg["ASP"] = np.where(agg["Units"]>0, agg["GMS"]/agg["Units"], np.nan)
    agg["CR"]  = np.where(agg["GV"]>0,    agg["Units"]/agg["GV"]*100, np.nan)
    agg["Month"] = pd.to_datetime(agg["year"].astype(str)+"-"+agg["month"].astype(str))
    agg = agg.sort_values("Month")

    for col in ["GMS","GV","Units","BA","ASP","CR"]:
        prev = agg[col].shift(1)
        agg[f"{col}_MoM"] = np.where(prev>0, (agg[col]-prev)/prev*100, np.nan)

    prev_y = agg[["year","month","GMS","GV","Units","BA","ASP","CR"]].copy()
    prev_y["year"] = prev_y["year"]+1
    merged = pd.merge(agg, prev_y, on=["year","month"], how="left", suffixes=("","_PY"))
    for col in ["GMS","GV","Units","BA","ASP","CR"]:
        agg[f"{col}_YoY"] = np.where(merged[f"{col}_PY"]>0, (agg[col]-merged[f"{col}_PY"])/merged[f"{col}_PY"]*100, np.nan)

    table_df = agg.sort_values("Month", ascending=False).copy()
    table_df["MonthStr"] = table_df["Month"].dt.strftime("%Y-%m")
    disp_cols = [
        "MonthStr","GMS","GV","Units","ASP","CR",
        "GMS_MoM","GV_MoM","Units_MoM","BA_MoM","ASP_MoM","CR_MoM",
        "GMS_YoY","GV_YoY","Units_YoY","BA_YoY","ASP_YoY","CR_YoY"
    ]
    table_df = table_df[disp_cols]

    _set_korean_font_if_possible()
    fig = plt.figure(figsize=(19.2, 10.8))
    fig.patch.set_facecolor("white")

    ax_wc = fig.add_axes([0.05, 0.60, 0.43, 0.32]); ax_wc.axis("off")
    ax_wc.set_title(f"ITK/BN Distribution by GMS  •  Month: {month_str}", fontsize=12, loc="left")
    if len(wc_labels) > 0:
        sizes = wc_vals / wc_vals.max()
        fs = 10 + (sizes**0.5) * 38
        cx, cy = 0.45, 0.45
        ax_wc.text(cx, cy, wc_labels[0], ha="center", va="center", fontsize=fs[0], color="#8AAAD6", alpha=0.85)
        rng = np.random.default_rng(42)
        xs = rng.uniform(0.05, 0.95, size=len(wc_labels)-1)
        ys = rng.uniform(0.05, 0.95, size=len(wc_labels)-1)
        for i, (tx, ty) in enumerate(zip(xs, ys), start=1):
            ax_wc.text(tx, ty, wc_labels[i], ha="center", va="center",
                       fontsize=fs[i], color="#AEC4E6", alpha=0.8, transform=ax_wc.transAxes)

    ax_d = fig.add_axes([0.56, 0.54, 0.38, 0.38])
    wedges, _ = ax_d.pie(donut_sizes, startangle=90, wedgeprops=dict(width=0.42, edgecolor="white"))
    ax_d.set_title("Sales% of ITK/BN", fontsize=12)
    ax_d.text(0, 0, _fmt_big(total_gms), ha="center", va="center", fontsize=18, fontweight="bold")
    ax_legend = fig.add_axes([0.80, 0.54, 0.18, 0.38]); ax_legend.axis("off")
    ax_legend.legend(wedges, donut_labels, loc="upper left", frameon=False, fontsize=9)

    ax_t = fig.add_axes([0.05, 0.06, 0.90, 0.44]); ax_t.axis("off")
    col_headers = ["Month","GMS","GV","Units","ASP","CR",
                   "GMS MoM","GV MoM","Units MoM","BA MoM","ASP MoM","CR MoM",
                   "GMS YoY","GV YoY","Units YoY","BA YoY","ASP YoY","CR YoY"]
    show = table_df.copy()
    show["GMS"]  = show["GMS"].map(lambda x: f"{x:,.0f}")
    show["GV"]   = show["GV"].map(lambda x: f"{x:,.0f}")
    show["Units"]= show["Units"].map(lambda x: f"{x:,.0f}")
    show["ASP"]  = show["ASP"].map(lambda x: f"{x:.2f}" if pd.notna(x) else "N/A")
    show["CR"]   = show["CR"].map(lambda x: fmt_pct(x))
    for c in ["GMS_MoM","GV_MoM","Units_MoM","BA_MoM","ASP_MoM","CR_MoM",
              "GMS_YoY","GV_YoY","Units_YoY","BA_YoY","ASP_YoY","CR_YoY"]:
        show[c] = show[c].map(lambda x: fmt_pct(x))

    cell_text = show.values.tolist()
    the_table = ax_t.table(cellText=cell_text, colLabels=col_headers, loc='center', cellLoc='center', colLoc='center')
    the_table.auto_set_font_size(False); the_table.set_fontsize(9); the_table.scale(1, 1.3)
    mom_idx = [col_headers.index(x) for x in ["GMS MoM","GV MoM","Units MoM","BA MoM","ASP MoM","CR MoM"]]
    yoy_idx = [col_headers.index(x) for x in ["GMS YoY","GV YoY","Units YoY","BA YoY","ASP YoY","CR YoY"]]
    for r in range(len(show)):
        for c_i, raw_col in zip(mom_idx, ["GMS_MoM","GV_MoM","Units_MoM","BA_MoM","ASP_MoM","CR_MoM"]):
            v = table_df.iloc[r][raw_col]; the_table[(r+1, c_i)].set_facecolor(_heat_color(v))
        for c_i, raw_col in zip(yoy_idx, ["GMS_YoY","GV_YoY","Units_YoY","BA_YoY","ASP_YoY","CR_YoY"]):
            v = table_df.iloc[r][raw_col]; the_table[(r+1, c_i)].set_facecolor(_heat_color(v))
    for c in range(len(col_headers)):
        the_table[(0, c)].set_facecolor("#F1F3F7"); the_table[(0, c)].set_text_props(weight='bold')

    fig.tight_layout()
    return _save_fig(fig, graph_name)

def create_graph12_top1_category_trends(df_asin, df_cid, graph_name="Graph 12"):
    def _get_col_safe(df, name):
        key = name.lower().replace(" ", "").replace("/", "")
        for c in df.columns:
            cc = str(c).lower().replace(" ", "").replace("/", "")
            if cc == key: return c
        raise KeyError(f"'{name}' 컬럼을 찾을 수 없습니다.")

    anchor = pd.to_datetime(df_cid["Month"], errors="coerce").max()
    if "Month" not in df_asin.columns:
        return None

    df_a = df_asin.copy()
    df_a["Month"] = pd.to_datetime(df_a["Month"], errors="coerce")
    df_a = df_a.dropna(subset=["Month"])

    col_itkbn = _get_col_safe(df_a, "ITK/BN")
    col_asin  = _get_col_safe(df_a, "ASIN")
    col_gms   = _get_col_safe(df_a, "GMS")
    col_gv    = _get_col_safe(df_a, "GV")
    col_units = _get_col_safe(df_a, "Units")

    cur = df_a[df_a["Month"].dt.to_period("M")==anchor.to_period("M")]
    if cur.empty: cur = df_a
    g_by = cur.groupby(col_itkbn, dropna=False)[col_gms].sum().sort_values(ascending=False)
    if g_by.empty:
        return None
    top1 = g_by.index[0]

    sub = df_a[df_a[col_itkbn]==top1].copy()
    sub["year"]  = sub["Month"].dt.year
    sub["month"] = sub["Month"].dt.month
    sub["BA_flag"] = (sub[col_units].fillna(0) > 0) | (sub[col_gms].fillna(0) > 0)
    agg = sub.groupby(["year","month"]).agg(
        GMS   =(col_gms,"sum"),
        GV    =(col_gv,"sum"),
        Units =(col_units,"sum"),
        BA    =(col_asin, lambda s: s[sub.loc[s.index,"BA_flag"]].nunique())
    ).reset_index()
    if agg.empty: return None
    agg["ASP"] = np.where(agg["Units"]>0, agg["GMS"]/agg["Units"], np.nan)
    agg["CR"]  = np.where(agg["GV"]>0,    agg["Units"]/agg["GV"]*100, np.nan)
    agg = finalize_year_month(agg, "year", "month")

    _set_korean_font_if_possible()
    fig, axes = plt.subplots(2,3, figsize=(19.2,10.8))
    fig.patch.set_facecolor("white")
    banner = fig.add_axes([0, 0.93, 1, 0.07])
    banner.set_facecolor("#0F6CBD"); banner.set_xticks([]); banner.set_yticks([])
    banner.text(0.5, 0.5, f"Top 1 카테고리 주요 지표 트렌드  •  {top1}",
                ha="center", va="center", color="white", fontsize=18, fontweight="bold")

    def draw_metric(ax, df, value_col, title, yfmt="raw", dec=1):
        _bi_theme(ax)
        years = sorted(df["year"].unique())
        hi_year = max(years)
        colors = {y: ("#F2B233" if y==hi_year else "#BDBDBD") for y in years}
        for y in years:
            suby = df[df["year"]==y].sort_values("month")
            x = suby["month"].to_numpy(); yv = suby[value_col].to_numpy()
            ax.plot(x, yv, marker="o", linewidth=2.2, color=colors[y], label=str(y))
            if y==hi_year and len(yv)>0 and pd.notna(yv[-1]):
                if yfmt=="k": txt = f"{yv[-1]/1000:.1f}K"
                elif yfmt=="pct": txt = f"{yv[-1]:.{dec}f}%"
                elif yfmt=="int": txt = f"{int(round(yv[-1]))}"
                else: txt = f"{yv[-1]:.{dec}f}"
                _label_last(ax, x, yv, txt, dy=8)
        ax.set_title(title, fontsize=14, color=PALETTE["dark"])
        ax.set_xticks(range(1,13)); ax.set_xlim(1,12)
        if yfmt=="k": ax.yaxis.set_major_formatter(_yfmt_k(dec))
        elif yfmt=="pct": ax.yaxis.set_major_formatter(_yfmt_decimal(dec, "%"))
        elif yfmt=="int": ax.yaxis.set_major_formatter(_yfmt_decimal(0, ""))        
        ax.legend(loc="upper left", fontsize=9, frameon=False)

    draw_metric(axes[0,0], agg, "GMS",   "매출",      yfmt="k",   dec=1)
    draw_metric(axes[0,1], agg, "GV",    "고객 유입", yfmt="k",   dec=1)
    draw_metric(axes[0,2], agg, "Units", "판매 수량", yfmt="k",   dec=1)
    draw_metric(axes[1,0], agg, "BA",    "판매 상품 개수", yfmt="int", dec=0)
    draw_metric(axes[1,1], agg, "ASP",   "판매 객단가", yfmt="raw", dec=2)
    draw_metric(axes[1,2], agg, "CR",    "구매 전환율", yfmt="pct", dec=1)

    fig.tight_layout(rect=[0,0,1,0.93])
    return _save_fig(fig, graph_name)



def create_graph13_top2_category_trends(df_asin, df_cid, graph_name="Graph 13"):
    def _get_col_safe(df, name):
        key = name.lower().replace(" ", "").replace("/", "")
        for c in df.columns:
            cc = str(c).lower().replace(" ", "").replace("/", "")
            if cc == key: return c
        raise KeyError(f"'{name}' 컬럼을 찾을 수 없습니다.")

    anchor = pd.to_datetime(df_cid["Month"], errors="coerce").max()
    if "Month" not in df_asin.columns:
        return None

    df_a = df_asin.copy()
    df_a["Month"] = pd.to_datetime(df_a["Month"], errors="coerce")
    df_a = df_a.dropna(subset=["Month"])

    col_itkbn = _get_col_safe(df_a, "ITK/BN")
    col_asin  = _get_col_safe(df_a, "ASIN")
    col_gms   = _get_col_safe(df_a, "GMS")
    col_gv    = _get_col_safe(df_a, "GV")
    col_units = _get_col_safe(df_a, "Units")

    cur = df_a[df_a["Month"].dt.to_period("M")==anchor.to_period("M")]
    if cur.empty: cur = df_a
    g_by = cur.groupby(col_itkbn, dropna=False)[col_gms].sum().sort_values(ascending=False)
    if len(g_by) < 2:
        return None
    top2 = g_by.index[1]
    sub = df_a[df_a[col_itkbn]==top2].copy()

    sub = df_a[df_a[col_itkbn]==top2].copy()
    sub["year"]  = sub["Month"].dt.year
    sub["month"] = sub["Month"].dt.month
    sub["BA_flag"] = (sub[col_units].fillna(0) > 0) | (sub[col_gms].fillna(0) > 0)
    agg = sub.groupby(["year","month"]).agg(
        GMS   =(col_gms,"sum"),
        GV    =(col_gv,"sum"),
        Units =(col_units,"sum"),
        BA    =(col_asin, lambda s: s[sub.loc[s.index,"BA_flag"]].nunique())
    ).reset_index()
    if agg.empty: return None
    agg["ASP"] = np.where(agg["Units"]>0, agg["GMS"]/agg["Units"], np.nan)
    agg["CR"]  = np.where(agg["GV"]>0,    agg["Units"]/agg["GV"]*100, np.nan)
    agg = finalize_year_month(agg, "year", "month")

    _set_korean_font_if_possible()
    fig, axes = plt.subplots(2,3, figsize=(19.2,10.8))
    fig.patch.set_facecolor("white")
    banner = fig.add_axes([0, 0.93, 1, 0.07])
    banner.set_facecolor("#0F6CBD"); banner.set_xticks([]); banner.set_yticks([])
    banner.text(0.5, 0.5, f"Top 2 카테고리 주요 지표 트렌드  •  {top2}",
                ha="center", va="center", color="white", fontsize=18, fontweight="bold")

    def draw_metric(ax, df, value_col, title, yfmt="raw", dec=1):
        _bi_theme(ax)
        years = sorted(df["year"].unique())
        hi_year = max(years)
        colors = {y: ("#F2B233" if y==hi_year else "#BDBDBD") for y in years}
        for y in years:
            suby = df[df["year"]==y].sort_values("month")
            x = suby["month"].to_numpy(); yv = suby[value_col].to_numpy()
            ax.plot(x, yv, marker="o", linewidth=2.2, color=colors[y], label=str(y))
            if y==hi_year and len(yv)>0 and pd.notna(yv[-1]):
                if yfmt=="k": txt = f"{yv[-1]/1000:.1f}K"
                elif yfmt=="pct": txt = f"{yv[-1]:.{dec}f}%"
                elif yfmt=="int": txt = f"{int(round(yv[-1]))}"
                else: txt = f"{yv[-1]:.{dec}f}"
                _label_last(ax, x, yv, txt, dy=8)
        ax.set_title(title, fontsize=14, color=PALETTE["dark"])
        ax.set_xticks(range(1,13)); ax.set_xlim(1,12)
        if yfmt=="k": ax.yaxis.set_major_formatter(_yfmt_k(dec))
        elif yfmt=="pct": ax.yaxis.set_major_formatter(_yfmt_decimal(dec, "%"))
        elif yfmt=="int": ax.yaxis.set_major_formatter(_yfmt_decimal(0, ""))        
        ax.legend(loc="upper left", fontsize=9, frameon=False)

    draw_metric(axes[0,0], agg, "GMS",   "매출",      yfmt="k",   dec=1)
    draw_metric(axes[0,1], agg, "GV",    "고객 유입", yfmt="k",   dec=1)
    draw_metric(axes[0,2], agg, "Units", "판매 수량", yfmt="k",   dec=1)
    draw_metric(axes[1,0], agg, "BA",    "판매 상품 개수", yfmt="int", dec=0)
    draw_metric(axes[1,1], agg, "ASP",   "판매 객단가", yfmt="raw", dec=2)
    draw_metric(axes[1,2], agg, "CR",    "구매 전환율", yfmt="pct", dec=1)

    fig.tight_layout(rect=[0,0,1,0.93])
    return _save_fig(fig, graph_name)
# =========================
# 그래프 16/17: Top1/Top2 Child ASIN 대시보드 (추가)
# =========================
def create_graph_asin_trends(df_asin, asin_value, rank_num=1, graph_name="Graph 16"):
    d = df_asin.copy()
    d["Month"] = pd.to_datetime(d["Month"], errors="coerce")
    d = d.dropna(subset=["Month"])
    asin_col  = _get_col_any(d, "ASIN", "Child ASIN", "ChildASIN")
    gms_col   = _get_col_any(d, "GMS")
    gv_col    = _get_col_any(d, "GV")
    units_col = _get_col_any(d, "Units")

    sub = d[d[asin_col] == asin_value].copy()
    if sub.empty:
        return None
    sub["year"]  = sub["Month"].dt.year
    sub["month"] = sub["Month"].dt.month

    agg = (sub.groupby(["year","month"])
              .agg(GMS=(gms_col,"sum"),
                   GV=(gv_col,"sum"),
                   Units=(units_col,"sum"))
              .reset_index())
    if agg.empty:
        return None

    agg["ASP"] = np.where(agg["Units"]>0, agg["GMS"]/agg["Units"], np.nan)
    agg["CR"]  = np.where(agg["GV"]>0,    agg["Units"]/agg["GV"]*100, np.nan)
    agg = finalize_year_month(agg, "year", "month")

    _set_korean_font_if_possible()
    fig, axes = plt.subplots(2,3, figsize=(19.2,10.8))
    fig.patch.set_facecolor("white")

    banner = fig.add_axes([0, 0.93, 1, 0.07])
    banner.set_facecolor("#0F6CBD"); banner.set_xticks([]); banner.set_yticks([])
    banner.text(0.5, 0.5, f"Top{rank_num} Child ASIN 주요 지표 트렌드  •  {asin_value}",
                ha="center", va="center", color="white", fontsize=18, fontweight="bold")

    def draw_metric(ax, dfm, col, title, yfmt="raw", dec=1):
        _bi_theme(ax)
        years = sorted(dfm["year"].unique())
        if not years: 
            return
        hi_year = max(years)
        colors = {y: ("#F2B233" if y==hi_year else "#BDBDBD") for y in years}
        for y in years:
            suby = dfm[dfm["year"]==y].sort_values("month")
            x = suby["month"].to_numpy(); yv = suby[col].to_numpy()
            ax.plot(x, yv, marker="o", linewidth=2.2, color=colors[y], label=str(y))
            if y==hi_year and len(yv)>0 and pd.notna(yv[-1]):
                if yfmt=="k":   txt = f"{yv[-1]/1000:.1f}K"
                elif yfmt=="pct": txt = f"{yv[-1]:.{dec}f}%"
                else:           txt = f"{yv[-1]:.{dec}f}"
                _label_last(ax, x, yv, txt, dy=8)
        ax.set_title(title, fontsize=14, color=PALETTE["dark"])
        ax.set_xticks(range(1,13)); ax.set_xlim(1,12)
        if yfmt=="k":   ax.yaxis.set_major_formatter(_yfmt_k(dec))
        elif yfmt=="pct": ax.yaxis.set_major_formatter(_yfmt_decimal(dec, "%"))
        else:           ax.yaxis.set_major_formatter(_yfmt_decimal(dec, ""))        
        ax.legend(loc="upper left", fontsize=9, frameon=False)

    draw_metric(axes[0,0], agg, "GMS",   "매출",        yfmt="k",   dec=1)
    draw_metric(axes[0,1], agg, "GV",    "고객 유입",   yfmt="k",   dec=1)
    draw_metric(axes[0,2], agg, "Units", "판매 수량",   yfmt="k",   dec=1)
    draw_metric(axes[1,0], agg, "ASP",   "평균 판매 가격", yfmt="raw", dec=2)
    draw_metric(axes[1,1], agg, "CR",    "구매 전환율", yfmt="pct", dec=1)
    axes[1,2].axis("off")

    fig.tight_layout(rect=[0,0,1,0.93])
    return _save_fig(fig, graph_name)

def build_top12_asin_graphs(df_asin, df_cid, top_n=2):
    try:
        anchor_dt = pd.to_datetime(df_cid["Month"], errors="coerce").max()
    except Exception:
        anchor_dt = pd.to_datetime(df_asin["Month"], errors="coerce").max()

    da = df_asin.copy()
    da["Month"] = pd.to_datetime(da["Month"], errors="coerce")
    da = da.dropna(subset=["Month"])
    asin_col = _get_col_any(da, "ASIN", "Child ASIN", "ChildASIN")
    gms_col  = _get_col_any(da, "GMS")

    cur = da[da["Month"].dt.to_period("M")==anchor_dt.to_period("M")]
    if cur.empty: cur = da

    top = (cur.groupby(asin_col, dropna=False)[gms_col]
             .sum()
             .sort_values(ascending=False)
             .head(top_n))
    asins = [str(x) for x in top.index.tolist()]
    marker_map = {}
    if len(asins) >= 1:
        p1 = create_graph_asin_trends(df_asin, asins[0], rank_num=1, graph_name="Graph 16")
        if p1: marker_map["그래프16"] = p1
    if len(asins) >= 2:
        p2 = create_graph_asin_trends(df_asin, asins[1], rank_num=2, graph_name="Graph 17")
        if p2: marker_map["그래프17"] = p2
    return marker_map

# =========================
# PPT 그림 삽입 유틸
# =========================
def _marker_key(text: str) -> str:
    s = (text or "").strip().lower()
    s = s.replace(" ", "")
    s = s.replace("graph", "그래프")
    m = re.search(r"그래프0*(\d+)", s, flags=re.I)
    return f"그래프{int(m.group(1))}" if m else ""

def _iter_all_shapes(shapes):
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in _iter_all_shapes(sh.shapes):
                yield sub
        else:
            yield sh

def _delete_shape(shape):
    try:
        el = shape._element
        el.getparent().remove(el)
    except Exception:
        pass

def insert_graphs_by_markers(ppt_path: str, marker_to_image: dict, save_as: str | None = None):
    prs = Presentation(ppt_path)
    placed = 0
    for slide in prs.slides:
        for sh in list(_iter_all_shapes(slide.shapes)):
            if getattr(sh, "has_text_frame", False) and sh.has_text_frame:
                text = "\n".join(p.text for p in sh.text_frame.paragraphs)
                key = _marker_key(text)
                img = marker_to_image.get(key)
                if img and os.path.exists(img):
                    left, top, width, height = sh.left, sh.top, sh.width, sh.height
                    _delete_shape(sh)
                    slide.shapes.add_picture(img, left, top, width=width, height=height)
                    placed += 1
    out = save_as or os.path.join(os.path.dirname(ppt_path), f"Updated_{os.path.basename(ppt_path)}")
    prs.save(out)
    return out, placed

# =========================
# 원본 process_graphs (업로드 파일 경로 기준)
# =========================
def process_graphs(cid_path, asin_path):
    ext = os.path.splitext(cid_path)[1].lower()
    df = pd.read_csv(cid_path) if ext == ".csv" else pd.read_excel(cid_path)
    col = lambda idx: (df.iloc[:, idx].tolist() if len(df.columns) > idx else [])

    dates = col(2)  # C열
    d = col(3); e = col(4); f = col(5); g = col(6); h = col(7); i = col(8); j = col(9)
    marker_to_path = {}
    marker_to_path["그래프1"]  = create_line_by_year(dates, d, "매출", "매출", "Graph 1", unit="K", decimal=1)
    marker_to_path["그래프2"]  = create_line_by_year(dates, e, "고객 유입", "GV", "Graph 2", unit="K", decimal=1)
    marker_to_path["그래프3"]  = create_line_by_year(dates, f, "판매 수량", "Units", "Graph 3", unit="K", decimal=1)
    marker_to_path["그래프4"]  = create_line_by_year(dates, g, "판매상품개수", "Buyable ASIN", "Graph 4", unit="none", decimal=0)
    marker_to_path["그래프5"]  = create_line_by_year(dates, h, "판매 객단가", "ASP", "Graph 5", unit="none", decimal=1)
    i_pct = [parse_number_any(x)*100 if parse_number_any(x) is not None else None for x in i]
    marker_to_path["그래프6"]  = create_line_by_year(dates, i_pct, "구매전환율", "Conversion %", "Graph 6", percentage=True)
    marker_to_path["그래프7"]  = create_line_by_year(dates, j, "SKU당 매출생산성", "GMS per BA", "Graph 7", unit="K", decimal=1)

    k = col(10); l = col(11)
    if g and (k or l):
        marker_to_path["그래프8"] = create_combo_ba_awagv_awas(dates, g, k, l, "Graph 8", "월별 BA / Discoverability / Conversion")
    cg = col(84); ch = col(85)
    if cg and ch:
        marker_to_path["그래프9"] = create_ipi_combo_graph(dates, cg, ch, "Graph 9", "IPI Score / Excess PCT")
    ak = col(36); al = col(37); am = col(38); an = col(39); ao = col(40)
    if d and (ak or al or am or an or ao):
        marker_to_path["그래프10"] = create_merchandising_graph(dates, d, ak, al, am, an, ao, "Graph 10", "Merchandising")

    df_cid = load_cid(cid_path)
    ext2 = os.path.splitext(asin_path)[1].lower()
    df_asin = pd.read_csv(asin_path) if ext2==".csv" else pd.read_excel(asin_path)

    marker_to_path["그래프11"] = create_graph11_itkbn_dashboard(df_asin, df_cid, None, "Graph 11")
    marker_to_path["그래프12"] = create_graph12_top1_category_trends(df_asin, df_cid, "Graph 12")
    marker_to_path["그래프13"] = create_graph13_top2_category_trends(df_asin, df_cid, "Graph 13")

    # ✅ Top1/Top2 ASIN 텍스트 마커 계산
    compute_topasin_metrics(df_asin, df_cid)
    # ✅ Top1/Top2 ASIN 대시보드(그래프16/17)
    try:
        top12_map = build_top12_asin_graphs(df_asin, df_cid, top_n=2)
        marker_to_path.update(top12_map)
    except Exception as _e:
        print(f"Top1/Top2 ASIN 그래프 생성 건너뜀: {_e}")

    az = col(51); bk = col(63); bv = col(74)
    if d and (az or bk or bv):
        marker_to_path["그래프18"] = create_ads_tacos_graph(dates, d, az, bk, bv, "Graph 18", "Ads Spend / TACOS")
    bc = col(54); bd = col(55); bn = col(65); bo = col(66); by = col(75); bz = col(76)
    if bc or bn or by:
        marker_to_path["그래프19"] = create_ads_impr_clicks_graph(dates, bc, bd, bn, bo, by, bz, "Graph 19", "Ads Impression / Clicks")
    be = col(56); bp = col(67); ca = col(77)
    if be or bp or ca:
        marker_to_path["그래프20"] = create_three_line_pct(dates, be, bp, ca, ['SP CTR','SB CTR','SD CTR'], 'CTR (%)', "Graph 20", "Ads CTR")
    bf = col(57); bq = col(68); cb = col(79)
    if bf or bq or cb:
        marker_to_path["그래프21"] = create_three_line_pct(dates, bf, bq, cb, ['SP CVR','SB CVR','SD CVR'], 'CVR (%)', "Graph 21", "Ads CVR")
    bg = col(58); br = col(69); cc = col(80)
    if bg or br or cc:
        marker_to_path["그래프22"] = create_three_line_pct_nolabel(dates, bg, br, cc, ['SP CPC','SB CPC','SD CPC'], 'CPC (%)', "Graph 22", "Ads CPC")
    bh = col(59); bs = col(70); cd = col(81)
    if bh or bs or cd:
        marker_to_path["그래프23"] = create_three_line_pct_nolabel(dates, bh, bs, cd, ['SP ACOS','SB ACOS','SD ACOS'], 'ACOS (%)', "Graph 23", "Ads ACOS")
    ba = col(52); bl = col(63); bw = col(74)
    if d and (ba or bl or bw):
        marker_to_path["그래프24"] = create_ads_sales_combo(dates, d, ba, bl, bw, "Graph 24", "Ads Sales")

    marker_to_path = {k:v for k,v in marker_to_path.items() if v}
    return marker_to_path, df_cid, df_asin

# =========================
# Streamlit 실행 버튼
# =========================
if st.button("🚀 PPT 생성하기", type="primary", disabled=not (cid_up and asin_up and ppt_up)):
    with st.spinner("임시 작업 폴더 준비 중..."):
        GRAPH_ROOT = tempfile.mkdtemp(prefix="mbr_")
        graphs_dir = ensure_graphs_folder()
        # 업로드 파일을 temp에 저장
        cid_path  = os.path.join(GRAPH_ROOT, cid_up.name)
        asin_path = os.path.join(GRAPH_ROOT, asin_up.name)
        ppt_path  = os.path.join(GRAPH_ROOT, ppt_up.name)
        with open(cid_path, "wb") as f: f.write(cid_up.getbuffer())
        with open(asin_path, "wb") as f: f.write(asin_up.getbuffer())
        with open(ppt_path, "wb") as f: f.write(ppt_up.getbuffer())

    try:
        with st.spinner("📈 그래프 생성 중 (1~12, 16~17, 18~24)..."):
            marker_to_path, df_cid, df_asin = process_graphs(cid_path, asin_path)
        st.success(f"그래프 생성 완료: {len(marker_to_path)}개")

        with st.spinner("🖼 PPT 템플릿에 그래프 자동 삽입 중..."):
            updated_ppt_path, placed = insert_graphs_by_markers(ppt_path, marker_to_path)
        with st.spinner("🔤 텍스트/표 마커 치환 + YTD 테이블 + Top10 ASIN 테이블 채우는 중..."):
            prs = Presentation(updated_ppt_path)
            process_ppt_markers(prs, df_cid)
            ytd_cnt = fill_ytd_table_on_slides(prs, df_cid)
        
            # ✅ Top10 Child ASIN 테이블 채우기
            try:
                fill_top10_asin_table(prs, df_asin, df_cid)
                st.success("Top10 Child ASIN 테이블 채움 완료")
            except Exception as e:
                st.error(f"Top10 ASIN 테이블 채우기 실패: {e}")
        
            final_path = os.path.join(GRAPH_ROOT, f"Filled_{os.path.basename(ppt_path)}")
            prs.save(final_path)
        
        st.success(f"완료! 그래프 삽입 {placed}개, YTD 테이블 {ytd_cnt}개 채움.")
        with open(final_path, "rb") as f:
            st.download_button(
                "📥 완성 PPT 다운로드",
                data=f.read(),
                file_name=os.path.basename(final_path),
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

        with st.expander("🔎 디버그(삽입된 그래프 목록)"):
            mp = pd.DataFrame(sorted(marker_to_path.items()), columns=["marker","image_path"])
            st.dataframe(mp, use_container_width=True)

    except Exception as e:
        st.error(f"에러 발생: {e}")
    finally:
        # 필요 시 작업폴더 정리하려면 아래 주석 해제
        # try: shutil.rmtree(GRAPH_ROOT); GRAPH_ROOT=None
        # except: pass
        pass

st.caption("Tip) 템플릿에는 '그래프1' 같은 정확한 텍스트 자리표시자를 넣으세요. 표/텍스트의 {marker}도 자동 치환됩니다.")
